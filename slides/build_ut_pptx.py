#!/usr/bin/env python3
"""Migrate slides/index.html content onto the official UT Austin 16:9 template."""

from __future__ import annotations

import copy
import re
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "05-23-V01-WSH_PowerPoint_16-9_TEMPLATE.pptx"
NOTES_PATH = ROOT / "talking-points.md"
OUT = ROOT / "UT_Speculative_Decoding.pptx"

ORANGE = RGBColor(0xBF, 0x57, 0x00)
ORANGE_DK = RGBColor(0xBE, 0x57, 0x00)
CHARCOAL = RGBColor(0x33, 0x3F, 0x48)
INK = RGBColor(0x16, 0x11, 0x0D)
MUTED = RGBColor(0x5B, 0x65, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LO = RGBColor(0x9B, 0x2C, 0x2C)
LINE = RGBColor(0xD6, 0xD2, 0xC3)
ROW_ALT = RGBColor(0xF7, 0xF4, 0xEF)

A_T = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"


def a_texts(shape):
    return list(shape._element.findall(f".//{A_T}"))


def set_a_texts(shape, values: list[str]):
    nodes = a_texts(shape)
    for i, node in enumerate(nodes):
        node.text = values[i] if i < len(values) else ""


def load_notes() -> dict[int, str]:
    text = NOTES_PATH.read_text()
    parts = re.split(r"\n## Slide (\d+)", text)
    notes: dict[int, str] = {}
    for i in range(1, len(parts), 2):
        num = int(parts[i])
        body = parts[i + 1]
        body = re.sub(r"^ — [^\n]+\n+", "", body)
        body = re.sub(r"\nTransition:.*", "", body, flags=re.S)
        notes[num] = body.strip()
    return notes


def set_run(run, text, *, size=18, bold=False, color=INK, name="Calibri", italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name


def clear_placeholder(slide, idx: int):
    try:
        ph = slide.placeholders[idx]
    except KeyError:
        return None
    el = ph._element
    el.getparent().remove(el)
    return None


def set_title(slide, title: str, *, size=24, kicker_text: str | None = None):
    ph = slide.placeholders[0]
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    set_run(run, title, size=size, bold=True, color=CHARCOAL, name="Calibri")
    if kicker_text:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        p2.space_after = Pt(0)
        run2 = p2.add_run()
        set_run(run2, kicker_text.upper(), size=10, bold=True, color=ORANGE, name="Calibri")


def add_notes(slide, text: str):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


def box(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0.75)):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        sh.shadow.inherit = False
    except Exception:
        pass
    sh.fill.solid()
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = line_w
    # remove default adjustment look
    spPr = sh._element.spPr
    for child in list(spPr):
        if child.tag.endswith("prstGeom"):
            child.set("prst", "rect")
    return sh


def tb(slide, l, t, w, h, text, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, name="Calibri", italic=False):
    sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color, name=name, italic=italic)
    return sh


def tb_lines(slide, l, t, w, h, lines, *, size=13, color=INK, bold=False, gap=6):
    sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0 if i == 0 else gap)
        p.space_after = Pt(0)
        run = p.add_run()
        set_run(run, line, size=size, bold=bold, color=color)
    return sh


def kicker(slide, text):
    # Section labels go in the title placeholder so they stay below the orange header.
    return text


def footer(slide, page, total=15):
    tb(slide, 0.50, 5.32, 6.8, 0.22,
       "Bharath Thiruveedula · Abhishek Kumar · Scalable Machine Learning",
       size=9, color=MUTED)
    tb(slide, 8.55, 5.32, 1.0, 0.22, f"{page:02d} / {total:02d}",
       size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def replace_a_t(shape, texts: list[str]):
    nodes = shape._element.xpath(".//a:t", namespaces=NSMAP_A)
    for i, node in enumerate(nodes):
        node.text = texts[i] if i < len(texts) else ""


def delete_slide(prs, index: int):
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def strip_section_ext(prs):
    P_EXT = "{http://schemas.openxmlformats.org/presentationml/2006/main}ext"
    for ext in list(prs._element.findall(f".//{P_EXT}")):
        uri = ext.get("uri")
        if uri in {
            "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}",
            "{EFAFB233-063F-42B5-8137-9DF3F51BA10A}",
        }:
            parent = ext.getparent()
            if parent is not None:
                parent.remove(ext)


def add_content_slide(prs):
    layout = prs.slide_masters[3].slide_layouts[1]  # white Title and Content
    slide = prs.slides.add_slide(layout)
    clear_placeholder(slide, 1)
    return slide


def style_table(table, header=True, hi_cells=None, lo_cells=None):
    hi_cells = hi_cells or set()
    lo_cells = lo_cells or set()
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            fill = cell.fill
            fill.solid()
            if r == 0 and header:
                fill.fore_color.rgb = ORANGE
                color = WHITE
                bold = True
                size = 10
            else:
                fill.fore_color.rgb = WHITE if r % 2 else ROW_ALT
                color = INK
                bold = c <= 1
                size = 11
                if (r, c) in hi_cells:
                    color = ORANGE
                    bold = True
                if (r, c) in lo_cells:
                    color = LO
                    bold = True
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c <= 1 else PP_ALIGN.RIGHT
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                for run in p.runs:
                    run.font.size = Pt(size)
                    run.font.bold = bold
                    run.font.color.rgb = color
                    run.font.name = "Calibri"
                    if c >= 2:
                        run.font.name = "Calibri"


def speed_bar(slide, l, t, w, h, speed, max_speed=2.5, under=False):
    track = box(slide, l, t, w, h, fill=LINE)
    frac = min(speed / max_speed, 1.0)
    if frac > 0:
        box(slide, l, t, w * frac, h, fill=LO if under else ORANGE)
    return track


def build():
    if not TEMPLATE.exists():
        raise SystemExit(f"missing template: {TEMPLATE}")
    notes = load_notes()
    tmp = ROOT / "_ut_build.pptx"
    shutil.copyfile(TEMPLATE, tmp)
    prs = Presentation(str(tmp))
    strip_section_ext(prs)

    # --- Slide 1: keep branded photo title, replace copy ---
    slide1 = prs.slides[0]
    for sh in slide1.shapes:
        if not sh.has_text_frame:
            continue
        raw = " ".join(p.text for p in sh.text_frame.paragraphs)
        if "Month" in raw or "20XX" in raw or "20xx" in raw:
            set_a_texts(sh, ["August ", "2026"])
        elif "Presenter" in raw or "speaker name" in raw:
            set_a_texts(sh, [
                "Bharath Thiruveedula · Abhishek Kumar",
                "Scalable Machine Learning,",
                " The University of Texas at Austin",
            ])
        elif "headline" in raw.lower():
            set_a_texts(sh, ["Speculative decoding", "under production", "serving stacks"])
            # Template headline is 48pt Arial Black all-caps; 42pt keeps this to 3 lines.
            for rPr in sh._element.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"):
                rPr.set("sz", "4200")
        elif "subtitle" in raw.lower() or "additional description" in raw.lower():
            set_a_texts(sh, [
                "Does the fastest draft method stay fastest when engine and load change?",
                "Project-reported range: 2.48× to 0.90×  ·  Llama-3.1-8B  ·  NVIDIA L40  ·  vLLM + SGLang",
            ])

    add_notes(slide1, notes.get(1, ""))

    # Delete sample slides 2-8 (keep branded title)
    for _ in range(7):
        delete_slide(prs, 1)

    # Helper to add remaining slides
    def new_slide():
        return add_content_slide(prs)

    # --- 02 Motivation ---
    s = new_slide(); set_title(s, "One token commits before the next step can start", kicker_text="Motivation  ·  the bottleneck")
    tb(s, 0.50, 1.72, 9.0, 0.36,
       "KV caching avoids recomputing the prompt, but it does not remove the token-to-token dependency.",
       size=15, color=CHARCOAL)
    steps = ["t1", "t2", "t3", "t4", "tn"]
    for i, tok in enumerate(steps):
        x = 0.55 + i * 1.85
        faint = i == 4
        box(s, x, 2.22, 1.45, 0.95, fill=WHITE, line=LINE if faint else ORANGE, line_w=Pt(1.5))
        tb(s, x, 2.28, 1.45, 0.42, "Forward pass", size=12, bold=True, color=MUTED if faint else CHARCOAL, align=PP_ALIGN.CENTER)
        tb(s, x, 2.68, 1.45, 0.36, "Target · 8B", size=11, color=MUTED, align=PP_ALIGN.CENTER)
        tb(s, x, 3.28, 1.45, 0.32, tok, size=16, bold=True, color=ORANGE if not faint else MUTED, align=PP_ALIGN.CENTER)
        if i < 4:
            tb(s, x + 1.38, 2.48, 0.50, 0.30, "→", size=16, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s, 0.50, 3.78, 9.0, 0.28, "Step i+1 cannot begin until token i is known. One full 8B target step per subsequent token.",
       size=12, italic=True, color=MUTED)
    box(s, 0.50, 4.18, 9.0, 0.90, fill=ROW_ALT, line=None)
    tb(s, 0.70, 4.32, 8.6, 0.62,
       "Speculation tries to amortize that expensive target-model step across several candidate tokens.",
       size=16, color=CHARCOAL, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 2); add_notes(s, notes.get(2, ""))

    # --- 03 Background ---
    s = new_slide(); set_title(s, "Draft several; verify together", kicker_text="Background  ·  how it works")
    cards = [
        ("01  ·  Draft", "Propose", "A cheaper mechanism predicts several candidate tokens ahead of the confirmed sequence."),
        ("02  ·  Target", "Verify", "The 8B model scores the candidate positions together in one verification step."),
        ("03  ·  Commit", "Resolve", "Accept the matching prefix; reject from the first mismatch and continue with the target prediction."),
    ]
    for i, (num, head, body) in enumerate(cards):
        x = 0.50 + i * 3.10
        box(s, x, 1.78, 2.90, 2.05, fill=WHITE, line=LINE)
        box(s, x, 1.78, 2.90, 0.08, fill=ORANGE)
        tb(s, x + 0.14, 1.92, 2.62, 0.26, num, size=11, bold=True, color=ORANGE)
        tb(s, x + 0.14, 2.20, 2.62, 0.40, head, size=22, bold=True, color=CHARCOAL)
        tb(s, x + 0.14, 2.62, 2.62, 1.05, body, size=13, color=CHARCOAL)
        if i < 2:
            tb(s, x + 2.78, 2.55, 0.32, 0.30, "→", size=18, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s, 0.50, 4.00, 4.5, 0.85,
       "Correctness.  With the exact acceptance procedure, the target model’s output distribution is preserved.",
       size=13, color=CHARCOAL)
    tb(s, 5.20, 4.00, 4.3, 0.85,
       "Trade-off.  The speedup must repay the cost of the drafter and the larger verification step.",
       size=13, color=CHARCOAL)
    tb(s, 0.50, 4.90, 9.0, 0.28, "Leviathan et al., ICML 2023", size=11, italic=True, color=MUTED)
    footer(s, 3); add_notes(s, notes.get(3, ""))

    # --- 04 Methods ---
    s = new_slide(); set_title(s, "The drafter changes; the target does not", kicker_text="Methods  ·  four mechanisms, one 8B target")
    methods = [
        ("AR baseline", "vLLM + SGLang", "No drafter. The 8B target generates one token at a time."),
        ("Classical 1B", "vLLM + SGLang", "A separate Llama-3.2-1B model drafts sequentially; the 8B target verifies."),
        ("Medusa", "vLLM only", "Additional heads predict future-token candidates in parallel; continuations use tree-based verification."),
        ("EAGLE-3", "vLLM + SGLang", "A specialized module performs direct token prediction using multi-layer target-model features."),
    ]
    for i, (name, meta, body) in enumerate(methods):
        col, row = i % 2, i // 2
        x, y = 0.50 + col * 4.70, 1.76 + row * 1.22
        box(s, x, y, 4.50, 1.10, fill=WHITE, line=LINE)
        box(s, x, y, 0.08, 1.10, fill=ORANGE)
        tb(s, x + 0.22, y + 0.08, 4.15, 0.32, name, size=18, bold=True, color=CHARCOAL)
        tb(s, x + 0.22, y + 0.40, 4.15, 0.22, meta.upper(), size=10, bold=True, color=ORANGE)
        tb(s, x + 0.22, y + 0.62, 4.15, 0.42, body, size=12, color=CHARCOAL)
    tb(s, 0.50, 4.28, 4.5, 0.70, "vLLM ablation: fixed draft lengths K = 1, 3, 5.", size=13, color=CHARCOAL)
    tb(s, 5.20, 4.28, 4.3, 0.70, "SGLang adaptive: a distinct acceptance-driven policy; not pooled with the vLLM grid.", size=13, color=CHARCOAL)
    tb(s, 0.50, 4.92, 9.0, 0.28,
       "Medusa is N/A on SGLang (not listed in 0.5.17 docs). EAGLE-3 checkpoints are engine-specific and were not mixed.",
       size=11, italic=True, color=MUTED)
    footer(s, 4); add_notes(s, notes.get(4, ""))

    # --- 05 Setup ---
    s = new_slide(); set_title(s, "One workload; two independent servers", kicker_text="Experimental setup  ·  measurement pipeline")
    nodes = [
        ("Input", "16 technical prompts", "Greedy  ·  T=0  ·  max 256 new tokens"),
        ("Shared target", "Llama-3.1-8B-Instruct", "BF16  ·  ctx 4096  ·  NVIDIA L40, 48 GB"),
        ("Serving stacks", "vLLM 0.27.1 + SGLang 0.5.17", "One process at a time  ·  c=1 and c=16"),
    ]
    for i, (lab, val, det) in enumerate(nodes):
        x = 0.50 + i * 3.10
        box(s, x, 1.76, 2.90, 1.45, fill=WHITE, line=LINE)
        box(s, x, 1.76, 2.90, 0.08, fill=ORANGE)
        tb(s, x + 0.14, 1.90, 2.62, 0.22, lab.upper(), size=10, bold=True, color=ORANGE)
        tb(s, x + 0.14, 2.14, 2.62, 0.50, val, size=16, bold=True, color=CHARCOAL)
        tb(s, x + 0.14, 2.66, 2.62, 0.42, det, size=12, color=MUTED)
        if i < 2:
            tb(s, x + 2.78, 2.28, 0.32, 0.30, "→", size=18, color=ORANGE, align=PP_ALIGN.CENTER)
    facts = [("Primary measure", "Output tokens/s"), ("Latency", "TTFT + TPOT"),
             ("Resources", "Peak VRAM"), ("Draft behavior", "vLLM accept counters")]
    for i, (k, v) in enumerate(facts):
        x = 0.50 + i * 2.32
        box(s, x, 3.36, 2.20, 0.70, fill=ROW_ALT)
        tb(s, x + 0.10, 3.40, 2.00, 0.22, k.upper(), size=9, bold=True, color=MUTED)
        tb(s, x + 0.10, 3.62, 2.00, 0.36, v, size=13, bold=True, color=CHARCOAL)
    box(s, 0.50, 4.18, 9.0, 0.92, fill=CHARCOAL)
    tb(s, 0.70, 4.28, 8.6, 0.72,
       "Known asymmetry: vLLM classical 1B and Medusa ran with --enforce-eager after CUDA-graph capture aborted; SGLang retained graphs. Every speedup is versus the same engine’s AR baseline.",
       size=13, color=WHITE)
    footer(s, 5); add_notes(s, notes.get(5, ""))

    # --- 06 Design ---
    s = new_slide(); set_title(s, "Fixed model; real serving-stack differences", kicker_text="Comparison design  ·  what the experiment can isolate")
    box(s, 0.50, 1.76, 4.45, 2.35, fill=WHITE, line=LINE)
    box(s, 0.50, 1.76, 4.45, 0.08, fill=ORANGE)
    tb(s, 0.70, 1.92, 4.05, 0.36, "Held fixed", size=20, bold=True, color=CHARCOAL)
    tb_lines(s, 0.70, 2.36, 4.05, 1.55, [
        "Target: Llama-3.1-8B-Instruct weights + tokenizer",
        "Workload: same prompts, decoding, GPU, and loads",
        "Measurement: shared client and metric definitions",
    ], size=14, color=CHARCOAL, gap=8)
    box(s, 5.15, 1.76, 4.35, 2.35, fill=WHITE, line=LINE)
    box(s, 5.15, 1.76, 4.35, 0.08, fill=ORANGE)
    tb(s, 5.35, 1.92, 4.00, 0.36, "Engine-specific", size=20, bold=True, color=CHARCOAL)
    tb_lines(s, 5.35, 2.36, 4.00, 1.55, [
        "Execution: scheduler, kernels, CUDA graphs, memory",
        "Coverage: which speculative methods each engine exposes",
        "Artifacts: legal EAGLE-3 checkpoint and configuration",
    ], size=14, color=CHARCOAL, gap=8)
    box(s, 0.50, 4.26, 9.0, 0.82, fill=CHARCOAL)
    tb(s, 0.70, 4.38, 8.6, 0.58,
       "This is an end-to-end serving comparison, not a controlled engine-only experiment. Similar rankings support, but do not prove, an algorithmic advantage.",
       size=14, color=WHITE)
    footer(s, 6); add_notes(s, notes.get(6, ""))

    # --- 07 Results c=1 ---
    s = new_slide(); set_title(s, "EAGLE-3 leads on both engines", kicker_text="Result  ·  concurrency 1  ·  project-reported  ·  NVIDIA L40")
    rows = [
        ("vLLM", "EAGLE-3  ·  K=5", 2.48, False),
        ("SGLang", "EAGLE-3  ·  adaptive", 2.36, False),
    ]
    for i, (eng, method, spd, under) in enumerate(rows):
        y = 1.78 + i * 1.05
        tb(s, 0.50, y, 1.55, 0.32, eng, size=18, bold=True, color=CHARCOAL)
        tb(s, 0.50, y + 0.32, 1.55, 0.28, method, size=11, color=MUTED)
        speed_bar(s, 2.20, y + 0.18, 5.55, 0.28, spd, under=under)
        tb(s, 2.20, y + 0.50, 5.55, 0.20, "0×          AR 1×                                      2.5×",
           size=9, color=MUTED)
        tb(s, 7.90, y, 1.60, 0.55, f"{spd:.2f}×", size=28, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
    box(s, 0.50, 3.95, 4.45, 1.05, fill=ROW_ALT)
    tb(s, 0.66, 4.05, 4.15, 0.85,
       "Closest cross-engine comparison: EAGLE-3 K=3 reports 2.22× on vLLM and 2.24× on SGLang.",
       size=13, color=CHARCOAL)
    box(s, 5.15, 3.95, 4.35, 1.05, fill=ROW_ALT)
    tb(s, 5.31, 4.05, 4.05, 0.85,
       "Engine-sensitive case: classical 1B reports 1.05× on eager-mode vLLM versus 1.90× on graph-enabled SGLang.",
       size=13, color=CHARCOAL)
    footer(s, 7); add_notes(s, notes.get(7, ""))

    # --- 08 Results c=16 ---
    s = new_slide(); set_title(s, "EAGLE still wins; one configuration regresses", kicker_text="Result  ·  concurrency 16  ·  project-reported  ·  NVIDIA L40")
    rows = [
        ("vLLM", "EAGLE-3  ·  K=5", 2.01, False),
        ("SGLang", "EAGLE-3  ·  adaptive", 1.49, False),
        ("vLLM", "Classical 1B  ·  eager*", 0.90, True),
    ]
    for i, (eng, method, spd, under) in enumerate(rows):
        y = 1.72 + i * 0.82
        tb(s, 0.50, y, 1.70, 0.28, eng, size=16, bold=True, color=CHARCOAL)
        tb(s, 0.50, y + 0.28, 1.70, 0.24, method, size=11, color=MUTED)
        speed_bar(s, 2.30, y + 0.16, 5.40, 0.24, spd, under=under)
        tb(s, 7.85, y, 1.65, 0.50, f"{spd:.2f}×", size=24, bold=True,
           color=LO if under else ORANGE, align=PP_ALIGN.RIGHT)
    tb(s, 0.50, 4.22, 9.0, 0.70,
       "The often-cited batch crossover did not occur for EAGLE at concurrency 16 on this L40; it did occur for eager-mode vLLM classical speculation. Throughput: 831.5 tok/s (vLLM K=5) and 628.6 tok/s (SGLang adaptive).",
       size=13, color=CHARCOAL)
    footer(s, 8); add_notes(s, notes.get(8, ""))

    # --- 09 Full matrix ---
    s = new_slide(); set_title(s, "The ranking is clear; the magnitude is stack-dependent", size=20, kicker_text="Result  ·  full matrix  ·  project-reported")
    tb(s, 0.50, 1.70, 9.0, 0.20,
       "tok/s = output throughput    ·    ×AR = same-engine speedup    ·    accept = vLLM mean accepted tokens/draft, c=1 / c=16",
       size=11, color=MUTED)

    def fill_table(slide, left, top, width, height, headers, data, hi=None, lo=None, col_widths=None):
        nrows, ncols = 1 + len(data), len(headers)
        shape = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top), Inches(width), Inches(height))
        tbl = shape.table
        if col_widths:
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Inches(w)
        for c, h in enumerate(headers):
            tbl.cell(0, c).text = h
        for r, row in enumerate(data, 1):
            for c, val in enumerate(row):
                tbl.cell(r, c).text = val
        style_table(tbl, hi_cells=hi or set(), lo_cells=lo or set())
        return tbl

    fill_table(
        s, 0.50, 1.94, 9.00, 1.72,
        ["vLLM 0.27", "c=1 tok/s", "×AR", "c=16 tok/s", "×AR", "accept"],
        [
            ["AR baseline", "40.2", "1.00", "413.7", "1.00", "—"],
            ["Classical 1B*", "42.3", "1.05", "373.6", "0.90", "2.56 / 2.69"],
            ["Medusa*", "55.9", "1.39", "463.7", "1.12", "0.96 / 0.94"],
            ["EAGLE-3 K=1", "66.8", "1.66", "600.3", "1.45", "0.76 / 0.75"],
            ["EAGLE-3 K=3", "89.1", "2.22", "741.9", "1.79", "1.65 / 1.62"],
            ["EAGLE-3 K=5", "99.5", "2.48", "831.5", "2.01", "2.05 / 2.07"],
        ],
        hi={(6, 1), (6, 2), (6, 3), (6, 4)},
        lo={(2, 4)},
        col_widths=[1.90, 1.30, 1.05, 1.40, 1.05, 2.30],
    )
    fill_table(
        s, 0.50, 3.72, 9.00, 1.42,
        ["SGLang 0.5", "c=1 tok/s", "×AR", "c=16 tok/s", "×AR", "accept"],
        [
            ["AR baseline", "35.1", "1.00", "421.0", "1.00", "—"],
            ["Classical 1B", "66.9", "1.90", "505.5", "1.20", "—"],
            ["Medusa", "N/A", "—", "N/A", "—", "not in 0.5.17"],
            ["EAGLE-3", "78.6", "2.24", "589.7", "1.40", "—"],
            ["EAGLE adaptive", "83.0", "2.36", "628.6", "1.49", "—"],
        ],
        hi={(5, 1), (5, 2), (5, 3), (5, 4)},
        col_widths=[1.90, 1.30, 1.05, 1.40, 1.05, 2.30],
    )
    tb(s, 0.50, 5.16, 9.0, 0.18,
       "*eager. Acceptance is vLLM-primary; blanks are unavailable, not zero.",
       size=10, italic=True, color=MUTED)
    footer(s, 9); add_notes(s, notes.get(9, ""))

    # --- 10 K ablation ---
    s = new_slide(); set_title(s, "Longer drafts helped through K=5 on this grid", kicker_text="Result  ·  EAGLE-3 fixed-K ablation  ·  vLLM")
    ks = [
        ("K=1", "0.76", "1.66×", "1.45×", False),
        ("K=3", "1.65", "2.22×", "1.79×", False),
        ("K=5", "2.05", "2.48×", "2.01×", True),
    ]
    for i, (k, acc, c1, c16, best) in enumerate(ks):
        x = 0.50 + i * 3.10
        box(s, x, 1.76, 2.90, 1.70, fill=WHITE, line=ORANGE if best else LINE, line_w=Pt(1.75 if best else 1.0))
        tb(s, x + 0.16, 1.86, 2.58, 0.40, k, size=24, bold=True, color=ORANGE if best else CHARCOAL)
        tb(s, x + 0.16, 2.32, 0.80, 0.50, acc, size=18, bold=True, color=CHARCOAL)
        tb(s, x + 0.16, 2.78, 0.80, 0.40, "accepted", size=10, color=MUTED)
        tb(s, x + 1.10, 2.32, 0.80, 0.50, c1, size=18, bold=True, color=CHARCOAL)
        tb(s, x + 1.10, 2.78, 0.80, 0.40, "c=1", size=10, color=MUTED)
        tb(s, x + 2.00, 2.32, 0.75, 0.50, c16, size=18, bold=True, color=CHARCOAL)
        tb(s, x + 2.00, 2.78, 0.75, 0.40, "c=16", size=10, color=MUTED)
    # simple grouped bars
    tb(s, 0.50, 3.58, 9.0, 0.24, "Speedup versus vLLM AR   ·   light = c=1, solid = c=16   ·   dashed line = AR baseline",
       size=11, color=MUTED)
    box(s, 0.50, 3.88, 9.0, 1.12, fill=WHITE, line=LINE)
    # baseline at 1.0 of 2.7 max => 1/2.7 of 8.2"
    maxs = 2.7
    base_x = 2.30
    bar_w = 6.70
    # AR line
    ar_off = (1.0 / maxs) * bar_w
    box(s, base_x + ar_off, 3.96, 0.02, 0.96, fill=MUTED)
    pairs = [(1.66, 1.45), (2.22, 1.79), (2.48, 2.01)]
    labels = ["K=1", "K=3", "K=5"]
    for i, ((c1, c16), lab) in enumerate(zip(pairs, labels)):
        y = 4.04 + i * 0.30
        tb(s, 0.62, y - 0.04, 1.50, 0.24, lab, size=11, color=CHARCOAL)
        box(s, base_x, y + 0.02, (c1 / maxs) * bar_w, 0.10, fill=RGBColor(0xD6, 0xA0, 0x7A))
        box(s, base_x, y + 0.12, (c16 / maxs) * bar_w, 0.10, fill=ORANGE)
    tb(s, 0.50, 5.04, 9.0, 0.22,
       "Fixed K values are separate vLLM configurations. SGLang adaptive is not a point on this curve. K>5 was not tested.",
       size=11, italic=True, color=MUTED)
    footer(s, 10); add_notes(s, notes.get(10, ""))

    # --- 11 Attribution ---
    s = new_slide(); set_title(s, "Acceptance helps; total cost decides", kicker_text="Why  ·  qualified interpretations")
    box(s, 0.50, 1.72, 9.0, 0.70, fill=ROW_ALT)
    tb(s, 0.70, 1.86, 8.6, 0.46,
       "speedup  =  accepted work  /  (draft cost + verification cost)",
       size=20, color=CHARCOAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cases = [
        ("Draft quality", "Medusa accepts 0.96 tokens/draft versus 2.05 for EAGLE K=5; their c=1 speedups are 1.39× and 2.48×."),
        ("Draft cost", "Classical 1B accepts more (2.56) yet reports only 1.05× on eager-mode vLLM. High acceptance alone is insufficient."),
        ("Serving stack", "The same 1B draft reports 1.05× on eager-mode vLLM and 1.90× on graph-enabled SGLang; this gap is confounded."),
        ("Offered load", "EAGLE remains above AR at c=16, but gains shrink as verification work grows with effective batch × K."),
    ]
    for i, (title, body) in enumerate(cases):
        col, row = i % 2, i // 2
        x, y = 0.50 + col * 4.70, 2.56 + row * 1.18
        box(s, x, y, 4.50, 1.08, fill=WHITE, line=LINE)
        tb(s, x + 0.16, y + 0.08, 4.18, 0.26, title.upper(), size=11, bold=True, color=ORANGE)
        tb(s, x + 0.16, y + 0.36, 4.18, 0.64, body, size=13, color=CHARCOAL)
    tb(s, 0.50, 4.96, 9.0, 0.28,
       "Consistent with those mechanisms; the experiment does not separately estimate each component’s causal contribution.",
       size=11, italic=True, color=MUTED)
    footer(s, 11); add_notes(s, notes.get(11, ""))

    # --- 12 Implications ---
    s = new_slide(); set_title(s, "The best observed choice depends on engine and load", kicker_text="Serving implications  ·  conditional on this setup")
    grid = [
        ["Load", "vLLM 0.27.1", "SGLang 0.5.17"],
        ["c = 1", "EAGLE-3, K=5\n99.5 tok/s  ·  2.48× AR", "EAGLE adaptive\n83.0 tok/s  ·  2.36× AR"],
        ["c = 16", "EAGLE-3, K=5\n831.5 tok/s  ·  2.01× AR", "EAGLE adaptive\n628.6 tok/s  ·  1.49× AR"],
    ]
    # header
    for i, h in enumerate(grid[0]):
        x = 0.50 + i * 3.10
        box(s, x, 1.76, 3.00, 0.46, fill=ORANGE)
        tb(s, x, 1.82, 3.00, 0.34, h.upper(), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for r, row in enumerate(grid[1:], 1):
        for c, val in enumerate(row):
            x = 0.50 + c * 3.10
            y = 2.22 + (r - 1) * 1.15
            box(s, x, y, 3.00, 1.10, fill=WHITE, line=LINE)
            if c == 0:
                tb(s, x, y + 0.32, 3.00, 0.46, val, size=22, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
            else:
                lines = val.split("\n")
                tb(s, x + 0.12, y + 0.16, 2.76, 0.40, lines[0], size=16, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
                tb(s, x + 0.12, y + 0.58, 2.76, 0.36, lines[1], size=12, color=MUTED, align=PP_ALIGN.CENTER)
    box(s, 0.50, 4.58, 9.0, 0.52, fill=CHARCOAL)
    tb(s, 0.66, 4.66, 8.7, 0.38,
       "Do not generalize the vLLM classical or Medusa rows: both include eager-mode execution. Medusa is N/A on SGLang.",
       size=13, color=WHITE)
    footer(s, 12); add_notes(s, notes.get(12, ""))

    # --- 13 Limitations ---
    s = new_slide(); set_title(s, "Read the trend, not the third decimal", kicker_text="Limitations  ·  scope of the evidence")
    box(s, 0.50, 1.76, 9.0, 0.70, fill=ROW_ALT)
    tb(s, 0.70, 1.88, 8.6, 0.48,
       "Outside this study: 70B or MoE models  ·  A100/H100  ·  sampled decoding  ·  production traffic  ·  confidence intervals",
       size=13, color=MUTED)
    box(s, 0.50, 2.58, 9.0, 1.20, fill=WHITE, line=ORANGE, line_w=Pt(1.75))
    tb(s, 0.70, 2.70, 8.6, 0.28, "What was measured", size=14, bold=True, color=ORANGE)
    tb(s, 0.70, 3.02, 8.6, 0.58,
       "Llama-3.1-8B-Instruct  ·  NVIDIA L40, 48 GB  ·  BF16  ·  greedy T=0  ·  16 short prompts  ·  vLLM 0.27.1 + SGLang 0.5.17",
       size=15, color=CHARCOAL)
    limits = [
        ("Statistical limit", "One published run per configuration after warmup; no error bars."),
        ("Internal confound", "vLLM 1B and Medusa used eager mode; SGLang retained graphs."),
        ("Asymmetric telemetry", "Acceptance analysis is vLLM-primary because equivalent counters were unavailable."),
    ]
    for i, (h, b) in enumerate(limits):
        x = 0.50 + i * 3.10
        box(s, x, 3.96, 2.90, 1.08, fill=WHITE, line=LINE)
        tb(s, x + 0.14, 4.06, 2.62, 0.28, h, size=13, bold=True, color=CHARCOAL)
        tb(s, x + 0.14, 4.36, 2.62, 0.58, b, size=12, color=CHARCOAL)
    footer(s, 13); add_notes(s, notes.get(13, ""))

    # --- 14 Conclusions ---
    s = new_slide(); set_title(s, "There is no universal speculative-decoding speedup", kicker_text="Conclusions  ·  answers to the research questions")
    tb(s, 0.50, 1.80, 2.40, 0.70, "2.48×", size=48, bold=True, color=ORANGE)
    tb(s, 0.50, 2.50, 2.40, 0.24, "BEST OBSERVED", size=11, bold=True, color=MUTED)
    tb(s, 3.00, 2.10, 0.70, 0.40, "TO", size=14, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    tb(s, 3.70, 1.80, 2.20, 0.70, "0.90×", size=48, bold=True, color=LO)
    tb(s, 3.70, 2.50, 2.20, 0.24, "WORST OBSERVED", size=11, bold=True, color=MUTED)
    points = [
        ("Ranking", "EAGLE-3 is first on both engines at c=1 and remains first at c=16."),
        ("Load", "EAGLE gains shrink but remain positive at c=16 on this L40."),
        ("Attribution", "Acceptance, draft cost, graph mode, checkpoint, and engine implementation all matter."),
    ]
    for i, (h, b) in enumerate(points):
        y = 2.90 + i * 0.52
        tb(s, 0.50, y, 9.0, 0.48, f"{h}:  {b}", size=15, color=CHARCOAL)
    tb(s, 0.50, 4.55, 9.0, 0.45,
       "A useful serving claim names the draft method, engine, hardware, and offered load.",
       size=16, italic=True, color=CHARCOAL)
    footer(s, 14); add_notes(s, notes.get(14, ""))

    # --- 15 Appendix ---
    s = new_slide(); set_title(s, "Evidence behind the deck", kicker_text="Appendix  ·  reproducibility and primary sources")
    box(s, 0.50, 1.76, 4.45, 2.85, fill=WHITE, line=LINE)
    tb(s, 0.70, 1.88, 4.05, 0.30, "PROJECT ARTIFACTS", size=12, bold=True, color=ORANGE)
    tb_lines(s, 0.70, 2.24, 4.05, 2.20, [
        "Repo: github.com/bharaththiruveedula/benchmark-speculative-decoding",
        "Measurements: results/table_latest.json",
        "Shared client: bench/run.py",
        "Server settings: configs/",
        "Coverage: 16/16 requests per published cell",
    ], size=13, color=CHARCOAL, gap=8)
    box(s, 5.15, 1.76, 4.35, 2.85, fill=WHITE, line=LINE)
    tb(s, 5.35, 1.88, 4.00, 0.30, "PRIMARY SOURCES", size=12, bold=True, color=ORANGE)
    tb_lines(s, 5.35, 2.24, 4.00, 2.20, [
        "Leviathan et al., ICML 2023 — speculative decoding",
        "Cai et al., 2024 — Medusa",
        "Li et al., 2025 — EAGLE-3",
        "vLLM and SGLang speculative-decoding docs",
        "NVIDIA L40 specifications",
    ], size=13, color=CHARCOAL, gap=8)
    url_box = tb(s, 0.50, 4.72, 9.0, 0.36,
       "https://github.com/bharaththiruveedula/benchmark-speculative-decoding  ·  Benchmark values are project-reported.",
       size=12, italic=False, color=ORANGE)
    url_box.text_frame.paragraphs[0].runs[0].hyperlink.address = (
        "https://github.com/bharaththiruveedula/benchmark-speculative-decoding"
    )
    footer(s, 15); add_notes(s, notes.get(15, ""))

    prs.core_properties.title = "Speculative decoding under production serving stacks"
    prs.core_properties.author = "Bharath Thiruveedula, Abhishek Kumar"
    prs.core_properties.subject = "Scalable Machine Learning — The University of Texas at Austin"
    prs.save(str(OUT))
    tmp.unlink(missing_ok=True)
    print(f"wrote {OUT}  slides={len(list(prs.slides._sldIdLst))}")


if __name__ == "__main__":
    build()

#!/usr/bin/env python3
"""Build the speculative-decoding deck using the official UT Austin template."""

from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "slides" / "05-23-V01-WSH_PowerPoint_16-9_TEMPLATE.pptx"
OUTPUT = ROOT / "slides" / "speculative-decoding-UT-template.pptx"
TALKING_POINTS = ROOT / "slides" / "talking-points.md"

UT_ORANGE = RGBColor(0xBF, 0x57, 0x00)
UT_DARK_ORANGE = RGBColor(0x9D, 0x3A, 0x00)
UT_CHARCOAL = RGBColor(0x33, 0x3F, 0x48)
UT_GRAY = RGBColor(0xD6, 0xD2, 0xC4)
UT_LIGHT_GRAY = RGBColor(0xF3, 0xF1, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x16, 0x11, 0x0D)
TEXT_GRAY = RGBColor(0x4A, 0x42, 0x3C)
MUTED = RGBColor(0x77, 0x70, 0x69)
RED = RGBColor(0x9B, 0x2C, 0x2C)
PALE_ORANGE = RGBColor(0xF5, 0xE2, 0xD3)

SLIDE_W = 10.0
SLIDE_H = 5.625
CONTENT_TOP = 1.37
CONTENT_BOTTOM = 5.28

with ZipFile(TEMPLATE) as _template_zip:
    UT_BACKGROUNDS = {
        "title": _template_zip.read("ppt/media/image1.jpg"),
        "dark": _template_zip.read("ppt/media/image2.jpg"),
        "gray": _template_zip.read("ppt/media/image3.jpg"),
        "white": _template_zip.read("ppt/media/image4.jpg"),
    }
    UT_TITLE_LOGO = _template_zip.read("ppt/media/image5.png")


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[index]


def remove_placeholders(slide) -> None:
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape._element.getparent().remove(shape._element)


def add_background(slide, theme: str) -> None:
    slide.shapes.add_picture(
        BytesIO(UT_BACKGROUNDS[theme]),
        0,
        0,
        width=Inches(SLIDE_W),
        height=Inches(SLIDE_H),
    )


def set_shape_text(
    shape,
    text: str,
    *,
    size: float,
    color: RGBColor,
    font: str = "Arial",
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 14,
    color: RGBColor = BLACK,
    font: str = "Arial",
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.text_frame.margin_left = Inches(margin)
    shape.text_frame.margin_right = Inches(margin)
    shape.text_frame.margin_top = Inches(margin)
    shape.text_frame.margin_bottom = Inches(margin)
    shape.text_frame.vertical_anchor = valign
    set_shape_text(
        shape,
        text,
        size=size,
        color=color,
        font=font,
        bold=bold,
        align=align,
    )
    return shape


def add_rule(slide, x: float, y: float, w: float, color=UT_ORANGE, width=1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x),
        Inches(y),
        Inches(x + w),
        Inches(y),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    radius: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_title(slide, title: str, theme: str, kicker: str | None = None) -> None:
    fg = WHITE if theme == "dark" else BLACK
    accent = RGBColor(0xF8, 0xA9, 0x40) if theme == "dark" else UT_ORANGE
    if kicker:
        add_text(
            slide,
            kicker.upper(),
            0.65,
            0.58,
            8.65,
            0.22,
            size=8,
            color=accent,
            font="Arial Black",
            bold=True,
        )
    add_text(
        slide,
        title.upper(),
        0.65,
        0.78,
        8.7,
        0.56,
        size=24,
        color=fg,
        font="Arial Black",
        bold=True,
    )


def add_footer(slide, number: int, theme: str) -> None:
    fg = WHITE if theme == "dark" else MUTED
    add_text(
        slide,
        "BHARATH THIRUVEEDULA · ABHISHEK KUMAR · SCALABLE MACHINE LEARNING",
        0.65,
        5.36,
        7.8,
        0.13,
        size=6.5,
        color=fg,
        font="Arial",
    )
    add_text(
        slide,
        f"{number:02d} / 15",
        8.75,
        5.34,
        0.6,
        0.15,
        size=7,
        color=fg,
        font="Arial Black",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_content_slide(prs: Presentation, theme: str, number: int, title: str, kicker: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(slide)
    add_background(slide, theme)
    add_title(slide, title, theme, kicker)
    add_footer(slide, number, theme)
    return slide


def add_label_value(
    slide,
    label: str,
    value: str,
    x: float,
    y: float,
    w: float,
    *,
    theme: str = "white",
    value_size: float = 16,
):
    fg = WHITE if theme == "dark" else BLACK
    muted = UT_GRAY if theme == "dark" else MUTED
    add_text(
        slide,
        label.upper(),
        x,
        y,
        w,
        0.18,
        size=7,
        color=muted,
        font="Arial Black",
        bold=True,
    )
    add_text(
        slide,
        value,
        x,
        y + 0.2,
        w,
        0.48,
        size=value_size,
        color=fg,
        font="Arial",
        bold=False,
    )


def add_bullet(
    slide,
    title: str,
    body: str,
    x: float,
    y: float,
    w: float,
    *,
    theme: str = "white",
    body_size: float = 11.5,
):
    fg = WHITE if theme == "dark" else BLACK
    secondary = UT_GRAY if theme == "dark" else TEXT_GRAY
    accent = RGBColor(0xF8, 0xA9, 0x40) if theme == "dark" else UT_ORANGE
    add_rect(slide, x, y + 0.07, 0.07, 0.07, fill=accent)
    add_text(
        slide,
        title,
        x + 0.16,
        y,
        w - 0.16,
        0.2,
        size=body_size,
        color=fg,
        font="Arial",
        bold=True,
    )
    add_text(
        slide,
        body,
        x + 0.16,
        y + 0.23,
        w - 0.16,
        0.5,
        size=body_size - 1,
        color=secondary,
        font="Arial",
    )


def add_metric(
    slide,
    value: str,
    label: str,
    x: float,
    y: float,
    w: float,
    *,
    theme: str,
    color: RGBColor | None = None,
    size: float = 38,
):
    fg = color or (RGBColor(0xF8, 0xA9, 0x40) if theme == "dark" else UT_ORANGE)
    muted = UT_GRAY if theme == "dark" else MUTED
    add_text(
        slide,
        value,
        x,
        y,
        w,
        0.62,
        size=size,
        color=fg,
        font="Arial Black",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        label.upper(),
        x,
        y + 0.62,
        w,
        0.25,
        size=7,
        color=muted,
        font="Arial Black",
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def parse_talking_points() -> dict[int, str]:
    text = TALKING_POINTS.read_text()
    matches = list(re.finditer(r"^## Slide (\d+) — .+$", text, flags=re.MULTILINE))
    sections: dict[int, str] = {}
    for idx, match in enumerate(matches):
        start = match.end()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        sections[int(match.group(1))] = text[start:end].strip()
    return sections


def add_notes(slide, text: str | None) -> None:
    if text:
        slide.notes_slide.notes_text_frame.text = text


def set_cell_text(cell, text: str, size: float, color: RGBColor, bold=False, align=PP_ALIGN.LEFT):
    cell.text = ""
    cell.margin_left = cell.margin_right = Inches(0.05)
    cell.margin_top = cell.margin_bottom = Inches(0.03)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_fill(cell, color: RGBColor) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    notes = parse_talking_points()

    # Use the official template's photographic title background and wordmark.
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_placeholders(title_slide)
    add_background(title_slide, "title")
    title_slide.shapes.add_picture(
        BytesIO(UT_TITLE_LOGO),
        Inches(7.63),
        Inches(0.35),
        width=Inches(2.05),
        height=Inches(1.0),
    )
    add_text(
        title_slide,
        "AUGUST 2026",
        0.6,
        0.5,
        2.0,
        0.22,
        size=10.5,
        color=WHITE,
        font="Arial Black",
        bold=True,
    )
    add_text(
        title_slide,
        "SPECULATIVE DECODING\nUNDER PRODUCTION\nSERVING STACKS",
        0.55,
        1.31,
        6.4,
        1.85,
        size=31,
        color=WHITE,
        font="Arial Black",
        bold=True,
    )
    add_rule(title_slide, 0.69, 3.4, 6.15, WHITE, 1.5)
    add_text(
        title_slide,
        "Project-reported speedups range from 2.48× to 0.90×\nLlama-3.1-8B · NVIDIA L40 · vLLM + SGLang · concurrency 1 + 16",
        0.6,
        3.65,
        6.3,
        0.55,
        size=12,
        color=WHITE,
        font="Arial",
    )
    add_text(
        title_slide,
        "BHARATH THIRUVEEDULA · ABHISHEK KUMAR\nScalable Machine Learning, The University of Texas at Austin",
        0.6,
        4.55,
        6.6,
        0.4,
        size=9.5,
        color=WHITE,
        font="Arial",
        bold=True,
    )
    add_notes(title_slide, notes.get(1))

    # 02 — Motivation
    slide = add_content_slide(
        prs,
        "white",
        2,
        "One token commits before the next step can start",
        "Motivation · the autoregressive bottleneck",
    )
    add_text(
        slide,
        "KV caching avoids recomputing the prompt, but it does not remove the token-to-token dependency.",
        0.65,
        CONTENT_TOP,
        8.65,
        0.35,
        size=13,
        color=TEXT_GRAY,
    )
    labels = ["8B FORWARD", "8B FORWARD", "8B FORWARD", "8B FORWARD"]
    tokens = ["t₁", "t₂", "t₃", "t₄"]
    for i, (label, token) in enumerate(zip(labels, tokens)):
        x = 0.72 + i * 2.18
        add_rect(slide, x, 2.05, 1.45, 0.65, fill=UT_LIGHT_GRAY, line=UT_GRAY, radius=True)
        add_text(slide, label, x, 2.22, 1.45, 0.2, size=10.5, color=BLACK, font="Arial Black", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, token, x + 0.5, 2.95, 0.45, 0.36, size=20, color=WHITE, font="Arial Black", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        token_shape = slide.shapes[-1]
        token_shape.fill.solid()
        token_shape.fill.fore_color.rgb = UT_ORANGE
        token_shape.line.color.rgb = UT_ORANGE
        if i < 3:
            add_text(slide, "→", x + 1.58, 2.2, 0.42, 0.3, size=20, color=UT_ORANGE, font="Arial Black", bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, "waits", x + 1.55, 2.65, 0.5, 0.2, size=7, color=MUTED, align=PP_ALIGN.CENTER)
    add_rule(slide, 0.65, 3.72, 8.65)
    add_text(
        slide,
        "SPECULATION TRIES TO AMORTIZE ONE EXPENSIVE TARGET STEP ACROSS SEVERAL CANDIDATE TOKENS.",
        0.65,
        3.9,
        8.65,
        0.55,
        size=17,
        color=UT_ORANGE,
        font="Arial Black",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "The prompt-prefill computation may already produce the first output token; the remaining generation steps are still serial.",
        1.2,
        4.62,
        7.6,
        0.3,
        size=9.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_notes(slide, notes.get(2))

    # 03 — Draft and verify
    slide = add_content_slide(
        prs, "gray", 3, "Draft several; verify together", "Background · speculative decoding"
    )
    steps = [
        ("01", "PROPOSE", "A cheaper drafter predicts several candidate tokens ahead."),
        ("02", "VERIFY", "The 8B target scores all candidate positions in one step."),
        ("03", "COMMIT", "Accept the matching prefix; correct at the first mismatch."),
    ]
    for i, (num, heading, body) in enumerate(steps):
        x = 0.65 + i * 3.0
        add_rect(slide, x, 1.55, 2.55, 1.55, fill=WHITE, line=UT_GRAY, radius=True)
        add_text(slide, num, x + 0.16, 1.7, 0.4, 0.25, size=11, color=UT_ORANGE, font="Arial Black", bold=True)
        add_text(slide, heading, x + 0.16, 2.0, 2.2, 0.3, size=17, color=BLACK, font="Arial Black", bold=True)
        add_text(slide, body, x + 0.16, 2.38, 2.2, 0.55, size=10.5, color=TEXT_GRAY)
        if i < 2:
            add_text(slide, "→", x + 2.58, 2.05, 0.38, 0.36, size=20, color=UT_ORANGE, font="Arial Black", bold=True, align=PP_ALIGN.CENTER)
    add_rule(slide, 0.65, 3.52, 8.65)
    add_bullet(slide, "Correctness", "Exact acceptance preserves the target model’s output distribution.", 0.75, 3.78, 4.0, body_size=11)
    add_bullet(slide, "Cost", "The accepted work must repay both drafting and verification.", 5.05, 3.78, 4.0, body_size=11)
    add_text(slide, "Leviathan et al., ICML 2023", 0.75, 4.83, 3.0, 0.18, size=7.5, color=MUTED)
    add_notes(slide, notes.get(3))

    # 04 — Methods
    slide = add_content_slide(
        prs, "white", 4, "The drafter changes; the target does not", "Methods · one shared 8B target"
    )
    methods = [
        ("AR BASELINE", "No drafter. The 8B target emits one token at a time.", "vLLM + SGLang"),
        ("CLASSICAL 1B", "A separate Llama-3.2-1B model drafts sequentially.", "vLLM + SGLang"),
        ("MEDUSA", "Additional heads predict future-token candidates in parallel.", "vLLM only"),
        ("EAGLE-3", "Direct token prediction from multi-layer target-model features.", "vLLM + SGLang"),
    ]
    for i, (name, mechanism, engines) in enumerate(methods):
        col, row = i % 2, i // 2
        x, y = 0.65 + col * 4.45, 1.48 + row * 1.45
        add_rule(slide, x, y, 4.05, UT_ORANGE if name == "EAGLE-3" else UT_GRAY, 2.2)
        add_text(slide, name, x, y + 0.14, 1.5, 0.28, size=14, color=BLACK, font="Arial Black", bold=True)
        add_text(slide, engines.upper(), x + 2.3, y + 0.16, 1.75, 0.2, size=7, color=UT_ORANGE, font="Arial Black", bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, mechanism, x, y + 0.5, 4.05, 0.46, size=10.5, color=TEXT_GRAY)
    add_rect(slide, 0.65, 4.38, 8.65, 0.58, fill=UT_LIGHT_GRAY, line=UT_GRAY)
    add_text(
        slide,
        "vLLM fixed-K ablation: K∈{1,3,5}   ·   SGLang adaptive: distinct acceptance-driven runtime policy",
        0.85,
        4.55,
        8.25,
        0.2,
        size=9.5,
        color=BLACK,
        font="Arial",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_notes(slide, notes.get(4))

    # 05 — Setup
    slide = add_content_slide(
        prs, "gray", 5, "One workload; two independent servers", "Experimental setup · measurement pipeline"
    )
    nodes = [
        ("INPUT", "16 technical prompts", "Greedy T=0 · max 256 tokens"),
        ("TARGET", "Llama-3.1-8B", "BF16 · context 4096 · L40"),
        ("SERVERS", "vLLM + SGLang", "0.27.1 · 0.5.17 · c=1,16"),
    ]
    for i, (label, value, detail) in enumerate(nodes):
        x = 0.65 + i * 3.0
        add_rect(slide, x, 1.45, 2.55, 1.14, fill=WHITE, line=UT_GRAY, radius=True)
        add_text(slide, label, x + 0.15, 1.6, 2.25, 0.18, size=7.5, color=UT_ORANGE, font="Arial Black", bold=True)
        add_text(slide, value, x + 0.15, 1.86, 2.25, 0.28, size=15, color=BLACK, font="Arial Black", bold=True)
        add_text(slide, detail, x + 0.15, 2.22, 2.25, 0.22, size=8.5, color=MUTED)
        if i < 2:
            add_text(slide, "→", x + 2.6, 1.82, 0.32, 0.35, size=19, color=UT_ORANGE, font="Arial Black", bold=True, align=PP_ALIGN.CENTER)
    facts = [("PRIMARY", "Output tok/s"), ("LATENCY", "TTFT + TPOT"), ("RESOURCE", "Peak VRAM"), ("DRAFT", "vLLM accept counters")]
    for i, (label, value) in enumerate(facts):
        x = 0.65 + i * 2.16
        add_label_value(slide, label, value, x, 2.95, 1.95, value_size=11.5)
    add_rect(slide, 0.65, 3.8, 8.65, 0.85, fill=UT_CHARCOAL, line=UT_CHARCOAL)
    add_text(slide, "KNOWN ASYMMETRY", 0.85, 3.98, 1.5, 0.2, size=8, color=RGBColor(0xF8, 0xA9, 0x40), font="Arial Black", bold=True)
    add_text(
        slide,
        "vLLM classical 1B and Medusa used --enforce-eager after CUDA-graph capture aborted; SGLang retained graphs.",
        2.15,
        3.94,
        6.9,
        0.36,
        size=10,
        color=WHITE,
        font="Arial",
        bold=True,
    )
    add_text(slide, "Every speedup is computed against the same engine’s AR baseline. All reported cells completed 16/16 requests.", 0.85, 4.39, 8.0, 0.18, size=7.5, color=UT_GRAY)
    add_notes(slide, notes.get(5))

    # 06 — Comparison design
    slide = add_content_slide(
        prs, "white", 6, "Fixed model; real serving-stack differences", "Comparison design · what can be isolated"
    )
    columns = [
        ("HELD FIXED", [
            "Target weights + tokenizer",
            "Prompts + decoding settings",
            "GPU + offered loads",
            "Shared client + metric definitions",
        ]),
        ("ENGINE-SPECIFIC", [
            "Scheduler + kernels",
            "CUDA-graph behavior",
            "Method support",
            "Legal EAGLE-3 checkpoint",
        ]),
    ]
    for col, (heading, items) in enumerate(columns):
        x = 0.65 + col * 4.45
        add_rule(slide, x, 1.48, 4.05, UT_ORANGE if col == 0 else UT_GRAY, 3)
        add_text(slide, heading, x, 1.67, 4.05, 0.28, size=15, color=BLACK, font="Arial Black", bold=True)
        for i, item in enumerate(items):
            add_bullet(slide, item, "", x, 2.12 + i * 0.5, 4.05, body_size=10.5)
    add_rect(slide, 0.65, 4.35, 8.65, 0.55, fill=UT_CHARCOAL, line=UT_CHARCOAL)
    add_text(
        slide,
        "END-TO-END SERVING COMPARISON — NOT A CONTROLLED ENGINE-ONLY EXPERIMENT",
        0.85,
        4.52,
        8.25,
        0.2,
        size=9,
        color=WHITE,
        font="Arial Black",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_notes(slide, notes.get(6))

    # 07 — c=1 results
    slide = add_content_slide(
        prs, "dark", 7, "EAGLE-3 leads on both engines", "Results · concurrency 1 · project-reported"
    )
    add_metric(slide, "2.48×", "vLLM · EAGLE-3 K=5 · 99.5 tok/s", 0.8, 1.55, 3.7, theme="dark", size=46)
    add_metric(slide, "2.36×", "SGLang · EAGLE adaptive · 83.0 tok/s", 5.0, 1.55, 3.7, theme="dark", size=46)
    add_rule(slide, 0.85, 2.73, 8.3, UT_GRAY)
    add_bullet(slide, "Closest comparison", "EAGLE-3 K=3: 2.22× on vLLM and 2.24× on SGLang.", 0.85, 3.0, 3.9, theme="dark", body_size=11)
    add_bullet(slide, "Engine-sensitive case", "Classical 1B: 1.05× on eager vLLM versus 1.90× on graph-enabled SGLang.", 5.05, 3.0, 3.95, theme="dark", body_size=11)
    add_text(slide, "Similar EAGLE rankings support, but do not prove, an algorithm-only advantage.", 1.0, 4.58, 8.0, 0.25, size=11, color=UT_GRAY, align=PP_ALIGN.CENTER)
    add_notes(slide, notes.get(7))

    # 08 — c=16 results
    slide = add_content_slide(
        prs, "dark", 8, "EAGLE still wins; one configuration regresses", "Results · concurrency 16 · project-reported"
    )
    metrics = [
        ("2.01×", "vLLM EAGLE K=5\n831.5 tok/s", RGBColor(0xF8, 0xA9, 0x40)),
        ("1.49×", "SGLang adaptive\n628.6 tok/s", RGBColor(0xF8, 0xA9, 0x40)),
        ("0.90×", "vLLM classical 1B*\n373.6 tok/s", RGBColor(0xFF, 0x9F, 0x9F)),
    ]
    for i, (value, label, color) in enumerate(metrics):
        x = 0.65 + i * 3.0
        add_rect(slide, x, 1.52, 2.55, 1.52, fill=RGBColor(0x2B, 0x35, 0x3D), line=UT_GRAY, radius=True)
        add_metric(slide, value, label, x + 0.08, 1.68, 2.39, theme="dark", color=color, size=37)
    add_rule(slide, 0.65, 3.43, 8.65, UT_GRAY)
    add_text(
        slide,
        "THE BATCH CROSSOVER DID NOT OCCUR FOR EAGLE AT c=16 ON THIS L40.",
        0.7,
        3.69,
        8.6,
        0.35,
        size=16,
        color=WHITE,
        font="Arial Black",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "*The vLLM classical row includes eager-mode execution; draft-model and eager overhead were not separately measured.",
        1.0,
        4.35,
        8.0,
        0.28,
        size=9,
        color=UT_GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_notes(slide, notes.get(8))

    # 09 — Full matrix
    slide = add_content_slide(
        prs, "white", 9, "The ranking is clear; the magnitude is stack-dependent", "Results · complete matrix"
    )
    rows = [
        ["ENGINE", "METHOD", "c=1 TOK/S", "×AR", "c=16 TOK/S", "×AR", "ACCEPT"],
        ["vLLM", "AR baseline", "40.2", "1.00", "413.7", "1.00", "—"],
        ["", "Classical 1B*", "42.3", "1.05", "373.6", "0.90", "2.56 / 2.69"],
        ["", "Medusa*", "55.9", "1.39", "463.7", "1.12", "0.96 / 0.94"],
        ["", "EAGLE-3 K=1", "66.8", "1.66", "600.3", "1.45", "0.76 / 0.75"],
        ["", "EAGLE-3 K=3", "89.1", "2.22", "741.9", "1.79", "1.65 / 1.62"],
        ["", "EAGLE-3 K=5", "99.5", "2.48", "831.5", "2.01", "2.05 / 2.07"],
        ["SGLang", "AR baseline", "35.1", "1.00", "421.0", "1.00", "—"],
        ["", "Classical 1B", "66.9", "1.90", "505.5", "1.20", "—"],
        ["", "Medusa", "N/A", "N/A", "N/A", "N/A", "—"],
        ["", "EAGLE-3", "78.6", "2.24", "589.7", "1.40", "—"],
        ["", "EAGLE adaptive", "83.0", "2.36", "628.6", "1.49", "—"],
    ]
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.65), Inches(1.43), Inches(8.65), Inches(3.6))
    table = shape.table
    widths = [0.85, 1.75, 1.1, 0.75, 1.15, 0.75, 1.25]
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        table.rows[r].height = Inches(0.29 if r else 0.32)
        for c, value in enumerate(row):
            header = r == 0
            fill = UT_CHARCOAL if header else (UT_LIGHT_GRAY if r in (2, 4, 6, 8, 10) else WHITE)
            if r == 7:
                fill = PALE_ORANGE
            set_cell_fill(table.cell(r, c), fill)
            text_color = WHITE if header else BLACK
            if (r == 6 and c in (2, 3, 4, 5)) or (r == 11 and c in (2, 3, 4, 5)):
                text_color = UT_ORANGE
            if r == 2 and c == 5:
                text_color = RED
            set_cell_text(
                table.cell(r, c),
                value,
                7.2 if not header else 7,
                text_color,
                bold=header or c in (0, 1),
                align=PP_ALIGN.RIGHT if c >= 2 else PP_ALIGN.LEFT,
            )
    add_text(
        slide,
        "*eager. Accept = c=1 / c=16 vLLM mean accepted tokens/draft. SGLang blanks are unavailable, not zero. Full metrics remain in results/table_latest.json.",
        0.65,
        5.06,
        8.65,
        0.2,
        size=6.8,
        color=MUTED,
    )
    add_notes(slide, notes.get(9))

    # 10 — K ablation
    slide = add_content_slide(
        prs, "gray", 10, "Longer drafts helped through K=5 on this grid", "Results · vLLM EAGLE-3 fixed-K ablation"
    )
    k_data = [
        ("K=1", "0.76", "1.66×", "1.45×"),
        ("K=3", "1.65", "2.22×", "1.79×"),
        ("K=5", "2.05", "2.48×", "2.01×"),
    ]
    for i, (k, accept, c1, c16) in enumerate(k_data):
        x = 0.65 + i * 3.0
        border = UT_ORANGE if i == 2 else UT_GRAY
        add_rect(slide, x, 1.5, 2.55, 2.45, fill=WHITE, line=border, radius=True)
        add_text(slide, k, x + 0.16, 1.7, 2.23, 0.35, size=20, color=UT_ORANGE if i == 2 else BLACK, font="Arial Black", bold=True)
        add_label_value(slide, "accepted / draft", accept, x + 0.16, 2.25, 1.0, value_size=18)
        add_label_value(slide, "c=1 speedup", c1, x + 1.25, 2.25, 1.0, value_size=18)
        add_label_value(slide, "c=16 speedup", c16, x + 0.16, 3.1, 2.08, value_size=18)
    add_text(slide, "ACCEPTANCE  0.76  →  1.65  →  2.05", 0.65, 4.27, 8.65, 0.3, size=13, color=UT_ORANGE, font="Arial Black", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "SGLang adaptive is a different policy and is not a point on this curve. K>5 was not tested.", 1.0, 4.72, 8.0, 0.25, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_notes(slide, notes.get(10))

    # 11 — Attribution
    slide = add_content_slide(
        prs, "white", 11, "Acceptance helps; total cost decides", "Discussion · qualified interpretation"
    )
    add_rect(slide, 1.05, 1.42, 7.9, 0.62, fill=UT_LIGHT_GRAY, line=UT_GRAY)
    add_text(slide, "speedup  ≈  accepted work  ÷  (draft cost + verification cost)", 1.2, 1.61, 7.6, 0.25, size=16, color=BLACK, font="Arial Black", bold=True, align=PP_ALIGN.CENTER)
    cases = [
        ("DRAFT QUALITY", "Medusa: accept 0.96, speedup 1.39×\nEAGLE K=5: accept 2.05, speedup 2.48×"),
        ("DRAFT COST", "Classical 1B accepts 2.56 yet reports only 1.05× on eager-mode vLLM."),
        ("SERVING STACK", "The same 1B draft reports 1.05× on vLLM and 1.90× on SGLang."),
        ("OFFERED LOAD", "EAGLE remains above AR at c=16, but its gains shrink."),
    ]
    for i, (heading, body) in enumerate(cases):
        col, row = i % 2, i // 2
        x, y = 0.65 + col * 4.45, 2.35 + row * 1.08
        add_rule(slide, x, y, 4.05, UT_ORANGE if i < 2 else UT_GRAY, 2)
        add_text(slide, heading, x, y + 0.12, 4.05, 0.2, size=9, color=UT_ORANGE, font="Arial Black", bold=True)
        add_text(slide, body, x, y + 0.38, 4.05, 0.48, size=9.5, color=TEXT_GRAY)
    add_text(slide, "These comparisons do not separately estimate each component’s causal contribution.", 1.0, 4.78, 8.0, 0.2, size=8, color=MUTED, align=PP_ALIGN.CENTER)

    # 12 — Implications
    slide = add_content_slide(
        prs, "white", 12, "The best observed choice depends on engine and load", "Serving implications · conditional on this setup"
    )
    data = [
        ("", "vLLM 0.27.1", "SGLang 0.5.17"),
        ("c=1", "EAGLE-3 K=5\n99.5 tok/s · 2.48×", "EAGLE adaptive\n83.0 tok/s · 2.36×"),
        ("c=16", "EAGLE-3 K=5\n831.5 tok/s · 2.01×", "EAGLE adaptive\n628.6 tok/s · 1.49×"),
    ]
    shape = slide.shapes.add_table(3, 3, Inches(0.65), Inches(1.55), Inches(8.65), Inches(2.55))
    table = shape.table
    table.columns[0].width = Inches(1.0)
    table.columns[1].width = table.columns[2].width = Inches(3.825)
    for r, row in enumerate(data):
        table.rows[r].height = Inches(0.48 if r == 0 else 1.03)
        for c, value in enumerate(row):
            fill = UT_CHARCOAL if r == 0 else (PALE_ORANGE if c == 1 else UT_LIGHT_GRAY)
            set_cell_fill(table.cell(r, c), fill)
            set_cell_text(
                table.cell(r, c),
                value,
                10 if r else 9,
                WHITE if r == 0 else BLACK,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
    add_rect(slide, 0.65, 4.38, 8.65, 0.52, fill=UT_CHARCOAL, line=UT_CHARCOAL)
    add_text(
        slide,
        "Do not generalize the eager-mode vLLM classical or Medusa rows.",
        0.85,
        4.54,
        8.25,
        0.2,
        size=9,
        color=WHITE,
        font="Arial Black",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 13 — Limitations
    slide = add_content_slide(
        prs, "gray", 13, "Read the trend, not the third decimal", "Limitations · scope of the evidence"
    )
    add_rect(slide, 0.65, 1.45, 8.65, 0.72, fill=UT_CHARCOAL, line=UT_CHARCOAL)
    add_text(
        slide,
        "OUTSIDE THIS STUDY: 70B / MoE · A100/H100 · sampled decoding · production traffic · confidence intervals",
        0.85,
        1.69,
        8.25,
        0.22,
        size=8.5,
        color=WHITE,
        font="Arial Black",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_rect(slide, 0.95, 2.38, 8.05, 1.05, fill=WHITE, line=UT_ORANGE)
    add_text(slide, "WHAT WAS MEASURED", 1.15, 2.58, 7.65, 0.22, size=9, color=UT_ORANGE, font="Arial Black", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Llama-3.1-8B · NVIDIA L40, 48 GB · BF16 · greedy T=0 · 16 prompts · vLLM 0.27.1 + SGLang 0.5.17", 1.15, 2.92, 7.65, 0.27, size=10.5, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    limits = [
        ("STATISTICS", "One published run per cell; no error bars."),
        ("CONFOUND", "vLLM 1B + Medusa eager; SGLang retained graphs."),
        ("TELEMETRY", "Acceptance analysis is vLLM-primary."),
    ]
    for i, (heading, body) in enumerate(limits):
        x = 0.65 + i * 3.0
        add_rule(slide, x, 3.85, 2.55, UT_ORANGE if i == 0 else UT_GRAY, 2)
        add_text(slide, heading, x, 4.02, 2.55, 0.2, size=8, color=UT_ORANGE, font="Arial Black", bold=True)
        add_text(slide, body, x, 4.3, 2.55, 0.45, size=9.5, color=TEXT_GRAY)

    # 14 — Conclusions
    slide = add_content_slide(
        prs, "dark", 14, "There is no universal speculative-decoding speedup", "Conclusions · answers to the research questions"
    )
    add_metric(slide, "2.48×", "best observed", 0.7, 1.52, 2.25, theme="dark", size=45)
    add_text(slide, "TO", 2.96, 1.88, 0.55, 0.25, size=8, color=UT_GRAY, font="Arial Black", bold=True, align=PP_ALIGN.CENTER)
    add_metric(slide, "0.90×", "worst observed", 3.48, 1.52, 2.25, theme="dark", color=RGBColor(0xFF, 0x9F, 0x9F), size=45)
    points = [
        ("RANKING", "EAGLE-3 ranks first on both engines at c=1 and c=16."),
        ("LOAD", "EAGLE gains shrink but remain positive at c=16 on this L40."),
        ("ATTRIBUTION", "Acceptance, draft cost, graph mode, checkpoint, and engine all matter."),
    ]
    for i, (heading, body) in enumerate(points):
        add_bullet(slide, heading, body, 6.05, 1.45 + i * 0.95, 3.15, theme="dark", body_size=10.5)
    add_rule(slide, 0.75, 4.32, 8.5, UT_GRAY)
    add_text(
        slide,
        "A USEFUL SERVING CLAIM NAMES THE DRAFT, ENGINE, HARDWARE, AND OFFERED LOAD.",
        0.8,
        4.55,
        8.4,
        0.3,
        size=13,
        color=WHITE,
        font="Arial Black",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 15 — Evidence
    slide = add_content_slide(
        prs, "white", 15, "Evidence behind the deck", "Appendix · reproducibility and primary sources"
    )
    add_rule(slide, 0.65, 1.46, 4.05, UT_ORANGE, 2)
    add_text(slide, "PROJECT ARTIFACTS", 0.65, 1.63, 4.05, 0.24, size=11, color=BLACK, font="Arial Black", bold=True)
    artifacts = [
        ("Measurements", "results/table_latest.json"),
        ("Shared client", "bench/run.py"),
        ("Server settings", "method-specific configs/ files"),
        ("Coverage", "16/16 completed requests per published cell"),
        ("Still needed", "raw logs, exact warmup trace, repeated runs"),
    ]
    for i, (heading, body) in enumerate(artifacts):
        add_bullet(slide, heading, body, 0.65, 2.02 + i * 0.55, 4.05, body_size=9)
    add_rule(slide, 5.1, 1.46, 4.2, UT_ORANGE, 2)
    add_text(slide, "PRIMARY SOURCES", 5.1, 1.63, 4.2, 0.24, size=11, color=BLACK, font="Arial Black", bold=True)
    sources = [
        "Leviathan et al. · Speculative decoding · ICML 2023",
        "Cai et al. · Medusa · 2024",
        "Li et al. · EAGLE-3 · 2025",
        "vLLM speculative-decoding documentation",
        "SGLang speculative-decoding documentation",
        "NVIDIA L40 specifications",
    ]
    for i, source in enumerate(sources):
        add_bullet(slide, source, "", 5.1, 2.02 + i * 0.48, 4.2, body_size=8.8)
    add_text(
        slide,
        "Benchmark values are project-reported; externally verifiable architecture and hardware statements use the sources above.",
        5.25,
        4.8,
        3.9,
        0.27,
        size=7.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build())
